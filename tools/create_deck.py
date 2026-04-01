"""
create_deck.py
--------------
Generates a professional PowerPoint slide deck from analyzed YouTube trend data.

Input:  .tmp/insights.json   (from analyze_trends.py)
        .tmp/charts/*.png    (from generate_charts.py)
Output: .tmp/ai_trend_report_YYYYMMDD.pptx
"""

import json
import os
import sys
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

INSIGHTS_FILE = ".tmp/insights.json"
CHARTS_DIR = ".tmp/charts"
OUTPUT_DIR = ".tmp"

# Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
C_BG_DARK = RGBColor(0x0F, 0x0F, 0x23)
C_BG_SURFACE = RGBColor(0x1A, 0x1A, 0x3E)
C_ACCENT_BLUE = RGBColor(0x4F, 0x8E, 0xF7)
C_ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
C_ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
C_ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xE8, 0xE8, 0xF0)
C_MUTED = RGBColor(0x90, 0x90, 0xB0)


def load_insights() -> dict:
    if not os.path.exists(INSIGHTS_FILE):
        print(f"ERROR: {INSIGHTS_FILE} not found. Run analyze_trends.py first.")
        sys.exit(1)
    with open(INSIGHTS_FILE, encoding="utf-8") as f:
        return json.load(f)


def chart_path(filename: str) -> str:
    return os.path.join(CHARTS_DIR, filename)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def blank_slide(prs: Presentation):
    """Add a completely blank slide."""
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def fill_slide_bg(slide, color: RGBColor):
    """Set slide background to solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text: str,
    left, top, width, height,
    font_size: int,
    bold: bool = False,
    color: RGBColor = C_LIGHT,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_image_if_exists(slide, path: str, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        # Placeholder box if chart wasn't generated
        add_rect(slide, left, top, width, height, C_BG_SURFACE)
        add_text(
            slide, f"[Chart not found:\n{os.path.basename(path)}]",
            left, top, width, height,
            font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER,
        )


def format_views(n) -> str:
    n = int(n)
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.0f}K"
    return str(n)


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_title(prs: Presentation, insights: dict):
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG_DARK)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, C_ACCENT_BLUE)

    # Gradient-ish surface panel on left half
    add_rect(slide, Inches(0.15), Inches(0), Inches(6.5), SLIDE_HEIGHT, C_BG_SURFACE)

    # Title text
    add_text(
        slide,
        "AI YouTube\nTrend Report",
        Inches(0.5), Inches(1.8),
        Inches(6), Inches(2.5),
        font_size=46, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT,
    )

    # Subtitle / date
    date_str = datetime.now().strftime("%B %Y")
    add_text(
        slide,
        f"{date_str}  ·  AI & Automation Niche",
        Inches(0.5), Inches(4.5),
        Inches(6), Inches(0.6),
        font_size=16, color=C_ACCENT_BLUE, align=PP_ALIGN.LEFT,
    )

    stats = insights["summary_stats"]
    add_text(
        slide,
        f"{insights['total_videos']:,} videos analyzed  ·  {stats['total_unique_channels']:,} channels",
        Inches(0.5), Inches(5.3),
        Inches(6), Inches(0.5),
        font_size=13, color=C_MUTED, align=PP_ALIGN.LEFT,
    )

    # Right side: accent decoration
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.3), Inches(0.06), C_ACCENT_BLUE)
    add_text(
        slide, "Powered by YouTube Data API v3",
        Inches(6.8), Inches(6.5), Inches(6), Inches(0.5),
        font_size=10, color=C_MUTED, align=PP_ALIGN.RIGHT,
    )


def slide_exec_summary(prs: Presentation, insights: dict):
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, C_ACCENT_GREEN)

    add_text(slide, "Executive Summary", Inches(0.3), Inches(0.25), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.95), Inches(12.7), Inches(0.04), C_ACCENT_GREEN)

    stats = insights["summary_stats"]
    recency = insights["recency_segments"]

    stat_cards = [
        ("Videos Analyzed", f"{insights['total_videos']:,}", C_ACCENT_BLUE),
        ("Unique Channels", f"{stats['total_unique_channels']:,}", C_ACCENT_GREEN),
        ("Avg Views", format_views(stats["avg_views"]), C_ACCENT_PURPLE),
        ("Avg Engagement", f"{stats['avg_engagement_rate']}%", C_ACCENT_ORANGE),
    ]

    card_w = Inches(2.9)
    card_h = Inches(1.6)
    card_top = Inches(1.2)
    gap = Inches(0.25)
    start_left = Inches(0.35)

    for i, (label, value, color) in enumerate(stat_cards):
        left = start_left + i * (card_w + gap)
        add_rect(slide, left, card_top, card_w, card_h, C_BG_SURFACE)
        add_rect(slide, left, card_top, card_w, Inches(0.07), color)
        add_text(slide, value, left, card_top + Inches(0.15), card_w, Inches(0.9),
                 font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, label, left, card_top + Inches(1.05), card_w, Inches(0.45),
                 font_size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Top video highlight
    if insights["top_videos_by_views"]:
        top_vid = insights["top_videos_by_views"][0]
        top_title = top_vid["title"][:80] + ("…" if len(top_vid["title"]) > 80 else "")
        add_rect(slide, Inches(0.35), Inches(3.05), Inches(12.6), Inches(1.35), C_BG_SURFACE)
        add_text(slide, "MOST VIEWED VIDEO", Inches(0.55), Inches(3.12), Inches(3), Inches(0.4),
                 font_size=9, bold=True, color=C_ACCENT_BLUE)
        add_text(slide, top_title, Inches(0.55), Inches(3.48), Inches(8), Inches(0.6),
                 font_size=14, bold=True, color=C_WHITE)
        add_text(slide,
                 f"{top_vid['channel_title']}  ·  {format_views(top_vid['view_count'])} views",
                 Inches(0.55), Inches(4.0), Inches(10), Inches(0.35),
                 font_size=11, color=C_MUTED)

    # Recency breakdown
    add_text(slide, "Content freshness:", Inches(0.35), Inches(4.6), Inches(4), Inches(0.4),
             font_size=12, bold=True, color=C_LIGHT)
    recency_text = (
        f"Last 7 days: {recency['last_7_days']} videos  ·  "
        f"Last 30 days: {recency['last_30_days']} videos  ·  "
        f"Last 90 days: {recency['last_90_days']} videos"
    )
    add_text(slide, recency_text, Inches(0.35), Inches(5.05), Inches(12.6), Inches(0.4),
             font_size=12, color=C_MUTED)


def slide_with_chart(
    prs: Presentation,
    title: str,
    subtitle: str,
    chart_file: str,
    accent_color: RGBColor,
    table_rows: list[dict] | None = None,
    table_headers: list[str] | None = None,
):
    """Generic slide: title + full-width chart (or chart + compact table)."""
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, accent_color)

    add_text(slide, title, Inches(0.3), Inches(0.2), Inches(12.8), Inches(0.65),
             font_size=26, bold=True, color=C_WHITE)
    add_text(slide, subtitle, Inches(0.3), Inches(0.85), Inches(12.8), Inches(0.4),
             font_size=12, color=C_MUTED)
    add_rect(slide, Inches(0.3), Inches(1.22), Inches(12.8), Inches(0.04), accent_color)

    if table_rows:
        # Chart on left, table on right
        chart_left = Inches(0.25)
        chart_top = Inches(1.35)
        chart_w = Inches(7.5)
        chart_h = Inches(5.85)
        add_image_if_exists(slide, chart_path(chart_file), chart_left, chart_top, chart_w, chart_h)

        # Compact table
        tbl_left = Inches(7.9)
        tbl_top = Inches(1.4)
        row_h = Inches(0.55)
        col_widths = [Inches(3.5), Inches(1.6)]

        if table_headers:
            for col_i, header in enumerate(table_headers):
                col_left = tbl_left + sum(col_widths[:col_i])
                add_rect(slide, col_left, tbl_top, col_widths[col_i], row_h, accent_color)
                add_text(slide, header, col_left + Inches(0.08), tbl_top + Inches(0.1),
                         col_widths[col_i] - Inches(0.1), row_h,
                         font_size=9, bold=True, color=C_WHITE)

        for row_i, row in enumerate(table_rows[:8]):
            row_top = tbl_top + row_h * (row_i + (1 if table_headers else 0))
            bg = C_BG_SURFACE if row_i % 2 == 0 else C_BG_DARK
            row_data = list(row.values())
            for col_i, cell_val in enumerate(row_data[:2]):
                col_left = tbl_left + sum(col_widths[:col_i])
                add_rect(slide, col_left, row_top, col_widths[col_i], row_h, bg)
                add_text(
                    slide, str(cell_val),
                    col_left + Inches(0.08), row_top + Inches(0.08),
                    col_widths[col_i] - Inches(0.1), row_h,
                    font_size=9, color=C_LIGHT,
                )
    else:
        # Full-width chart
        add_image_if_exists(
            slide, chart_path(chart_file),
            Inches(0.25), Inches(1.35),
            Inches(12.85), Inches(5.9),
        )


def slide_recommendations(prs: Presentation, insights: dict):
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, C_ACCENT_ORANGE)

    add_text(slide, "Content Recommendations", Inches(0.3), Inches(0.2), Inches(12.8), Inches(0.65),
             font_size=26, bold=True, color=C_WHITE)
    add_text(slide, "Data-driven insights from analyzed AI niche videos",
             Inches(0.3), Inches(0.85), Inches(12.8), Inches(0.4),
             font_size=12, color=C_MUTED)
    add_rect(slide, Inches(0.3), Inches(1.22), Inches(12.8), Inches(0.04), C_ACCENT_ORANGE)

    recs = insights.get("content_recommendations", [])
    if not recs:
        recs = ["Run the analysis pipeline to generate data-driven recommendations."]

    icons = ["01", "02", "03", "04", "05"]
    for i, rec in enumerate(recs[:5]):
        top = Inches(1.45) + i * Inches(1.1)
        # Number badge
        add_rect(slide, Inches(0.3), top, Inches(0.45), Inches(0.7), C_ACCENT_ORANGE)
        add_text(slide, icons[i], Inches(0.3), top, Inches(0.45), Inches(0.7),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Rec card
        add_rect(slide, Inches(0.82), top, Inches(12.3), Inches(0.9), C_BG_SURFACE)
        add_text(slide, rec, Inches(0.97), top + Inches(0.12), Inches(12.0), Inches(0.75),
                 font_size=13, color=C_LIGHT)


def slide_methodology(prs: Presentation, insights: dict):
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, C_MUTED)

    add_text(slide, "Methodology", Inches(0.3), Inches(0.2), Inches(12.8), Inches(0.65),
             font_size=26, bold=True, color=C_WHITE)
    add_rect(slide, Inches(0.3), Inches(0.9), Inches(12.8), Inches(0.04), C_MUTED)

    date_str = insights.get("generated_at", "")[:10]
    lines = [
        f"Data source:    YouTube Data API v3",
        f"Date collected: {date_str}",
        f"Time window:    Last 90 days",
        f"Total videos:   {insights['total_videos']:,}",
        f"Search keywords used:",
    ]
    keywords_used = [
        "AI automation 2025", "AI agents tutorial", "artificial intelligence tools",
        "LLM workflow automation", "ChatGPT automation", "AI productivity tools",
        "machine learning tutorial", "no code AI automation",
    ]

    top = Inches(1.1)
    for line in lines:
        add_text(slide, line, Inches(0.5), top, Inches(12), Inches(0.45),
                 font_size=13, color=C_LIGHT)
        top += Inches(0.45)

    top += Inches(0.1)
    for kw in keywords_used:
        add_rect(slide, Inches(0.5), top, Inches(0.08), Inches(0.32), C_ACCENT_BLUE)
        add_text(slide, kw, Inches(0.72), top, Inches(11), Inches(0.4),
                 font_size=12, color=C_MUTED)
        top += Inches(0.38)

    add_text(slide, "Generated automatically by the WAT YouTube Trend Analysis pipeline.",
             Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             font_size=10, color=C_MUTED)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    insights = load_insights()
    prs = new_presentation()

    print("Building slide deck...")

    # Slide 1: Title
    slide_title(prs, insights)
    print("  Slide 1: Title")

    # Slide 2: Executive Summary
    slide_exec_summary(prs, insights)
    print("  Slide 2: Executive Summary")

    # Slide 3: Top Videos by Views
    top_views_table = [
        {
            "title": v["title"][:40] + "…" if len(v["title"]) > 40 else v["title"],
            "views": format_views(v["view_count"]),
        }
        for v in insights["top_videos_by_views"][:8]
    ]
    slide_with_chart(
        prs,
        title="Top Videos by Views",
        subtitle="Highest viewed AI/automation videos published in the last 90 days",
        chart_file="top_videos_views.png",
        accent_color=C_ACCENT_BLUE,
        table_rows=top_views_table,
        table_headers=["Title", "Views"],
    )
    print("  Slide 3: Top Videos by Views")

    # Slide 4: Top Videos by Engagement
    top_eng_table = [
        {
            "title": v["title"][:40] + "…" if len(v["title"]) > 40 else v["title"],
            "engagement": f"{v['engagement_rate']}%",
        }
        for v in insights["top_videos_by_engagement"][:8]
    ]
    slide_with_chart(
        prs,
        title="Top Videos by Engagement Rate",
        subtitle="Engagement = (likes + comments) / views × 100 — measures audience resonance",
        chart_file="top_videos_engagement.png",
        accent_color=C_ACCENT_GREEN,
        table_rows=top_eng_table,
        table_headers=["Title", "Eng. Rate"],
    )
    print("  Slide 4: Top Videos by Engagement")

    # Slide 5: Trending Keywords
    slide_with_chart(
        prs,
        title="Trending Keywords in AI Niche Titles",
        subtitle="Most frequent meaningful words appearing in top-performing video titles",
        chart_file="trending_keywords.png",
        accent_color=C_ACCENT_PURPLE,
    )
    print("  Slide 5: Trending Keywords")

    # Slide 6: Channel Leaderboard
    ch_table = [
        {
            "channel": c["channel_title"][:35] + "…" if len(c["channel_title"]) > 35 else c["channel_title"],
            "total_views": format_views(c["total_views"]),
        }
        for c in insights["channel_leaderboard"][:8]
    ]
    slide_with_chart(
        prs,
        title="Top Channels by Total Views",
        subtitle="Channels with highest combined views across analyzed videos in this niche",
        chart_file="channel_leaderboard.png",
        accent_color=C_ACCENT_PURPLE,
        table_rows=ch_table,
        table_headers=["Channel", "Total Views"],
    )
    print("  Slide 6: Channel Leaderboard")

    # Slide 7: Upload Activity
    slide_with_chart(
        prs,
        title="Upload Activity Over Last 90 Days",
        subtitle="Number of new AI/automation videos published per week — shows niche momentum",
        chart_file="upload_trend.png",
        accent_color=C_ACCENT_ORANGE,
    )
    print("  Slide 7: Upload Activity")

    # Slide 8: Content Recommendations
    slide_recommendations(prs, insights)
    print("  Slide 8: Content Recommendations")

    # Slide 9: Methodology
    slide_methodology(prs, insights)
    print("  Slide 9: Methodology")

    # Save
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    date_str = datetime.now().strftime("%Y%m%d")
    output_path = os.path.join(OUTPUT_DIR, f"ai_trend_report_{date_str}.pptx")
    prs.save(output_path)
    print(f"\nDone. Deck saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
